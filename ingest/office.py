from __future__ import annotations

import os
import re
from typing import Any, Dict, List, Optional, Tuple

from .base import BaseParser
from ..core.types import ParsedDocument


class OfficeParser(BaseParser):

    OFFICE_FORMATS = {"docx", "pptx", "xlsx", "doc", "ppt", "xls", "odt", "odp", "ods"}

    def supports(self, format: str) -> bool:
        return format.lower() in self.OFFICE_FORMATS

    def parse(self, source: str) -> ParsedDocument:
        if not self._file_exists(source):
            raise FileNotFoundError(f"File not found: {source}")

        ext = os.path.splitext(source)[1].lower().lstrip(".")

        dispatch = {
            "docx": self._parse_docx,
            "pptx": self._parse_pptx,
            "xlsx": self._parse_xlsx,
        }

        handler = dispatch.get(ext)
        if handler:
            return handler(source)

        if ext in ("doc", "ppt", "xls"):
            raise ImportError(
                f"Legacy .{ext} format is not supported directly. "
                f"Please convert to .{ext}x format first. "
                f"You can do this by opening in your office application and "
                f"saving as .{ext}x."
            )

        if ext in ("odt", "odp", "ods"):
            raise ImportError(
                f"OpenDocument .{ext} format requires additional libraries. "
                f"Convert to .{ext.replace('o', 'x').replace('t', 'x').replace('d', 'x')}x or "
                f"install odfpy: pip install odfpy"
            )

        raise ValueError(f"Unsupported office format: .{ext}")

    def _parse_docx(self, path: str) -> ParsedDocument:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX parsing. "
                "Install with: pip install python-docx"
            )

        doc = DocxDocument(path)
        pages: List[Tuple[int, str]] = []
        sections: List[str] = []
        current_section = ""
        page_num = 1

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = (para.style.name or "").lower()

            if "heading" in style_name:
                level = ""
                for ch in style_name:
                    if ch.isdigit():
                        level += ch
                prefix = "#" * int(level) if level else "##"
                current_section = text
                sections.append(f"{prefix} {text}\n")
                continue

            if text:
                sections.append(text + "\n")

        for table in doc.tables:
            table_text = self._extract_docx_table(table)
            if table_text:
                sections.append(table_text + "\n")

        full_text = "\n".join(sections)
        full_text = self._strip_control_chars(full_text)
        full_text = self._normalize_whitespace(full_text)

        pages.append((page_num, full_text))

        metadata = {
            "format": "docx",
            "file_name": os.path.basename(path),
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        }

        try:
            core = doc.core_properties
            if core.title:
                metadata["title"] = core.title
            if core.author:
                metadata["author"] = core.author
            if core.subject:
                metadata["subject"] = core.subject
        except Exception:
            pass

        return self._make_result(full_text, pages, metadata)

    def _extract_docx_table(self, table: Any) -> str:
        rows: List[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        return "\n".join(rows) if rows else ""

    def _parse_pptx(self, path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. "
                "Install with: pip install python-pptx"
            )

        prs = Presentation(path)
        pages: List[Tuple[int, str]] = []
        all_text_parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: List[str] = []

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"[Speaker Notes] {notes_text}")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            is_title = False
                            if para.font and para.font.size:
                                if para.font.size >= 200000:
                                    is_title = True
                            if shape.has_text_frame and shape == slide.shapes.title if hasattr(slide.shapes, "title") else False:
                                is_title = True
                            if is_title:
                                slide_parts.append(f"# {text}")
                            else:
                                slide_parts.append(text)

                if shape.has_table:
                    table_text = self._extract_pptx_table(shape.table)
                    if table_text:
                        slide_parts.append(table_text)

            if shape.shape_type == 13 and hasattr(shape, "image"):
                try:
                    slide_parts.append(f"[Image: {shape.image.content_type}]")
                except Exception:
                    slide_parts.append("[Image]")

            slide_text = "\n".join(slide_parts)
            if slide_text.strip():
                pages.append((slide_num, slide_text))
                all_text_parts.append(f"--- Slide {slide_num} ---\n{slide_text}")

        full_text = "\n\n".join(all_text_parts)
        full_text = self._strip_control_chars(full_text)

        metadata = {
            "format": "pptx",
            "file_name": os.path.basename(path),
            "slide_count": len(prs.slides),
        }

        return self._make_result(full_text, pages, metadata)

    def _extract_pptx_table(self, table: Any) -> str:
        rows: List[str] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            if any(cells):
                rows.append(" | ".join(cells))
        return "\n".join(rows) if rows else ""

    def _parse_xlsx(self, path: str) -> ParsedDocument:
        try:
            import openpyxl
        except ImportError:
            raise ImportError(
                "openpyxl is required for XLSX parsing. "
                "Install with: pip install openpyxl"
            )

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        pages: List[Tuple[int, str]] = []
        all_text_parts: List[str] = []

        for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
            ws = wb[sheet_name]
            sheet_lines: List[str] = [f"Sheet: {sheet_name}\n"]
            headers: Optional[List[str]] = None
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                if all(cell is None for cell in row):
                    continue

                str_cells = []
                for cell in row:
                    if cell is None:
                        str_cells.append("")
                    else:
                        str_cells.append(str(cell).strip())

                if headers is None and any(str_cells):
                    headers = str_cells
                    header_line = " | ".join(str_cells)
                    sheet_lines.append(header_line)
                    sheet_lines.append("-" * len(header_line))
                    continue

                if any(str_cells):
                    if headers and len(str_cells) == len(headers):
                        parts = []
                        for header, val in zip(headers, str_cells):
                            if val:
                                parts.append(f"{header}: {val}")
                        if parts:
                            sheet_lines.append(", ".join(parts))
                    else:
                        sheet_lines.append(" | ".join(str_cells))
                    row_count += 1

            if row_count > 0:
                sheet_text = "\n".join(sheet_lines)
                pages.append((sheet_idx, sheet_text))
                all_text_parts.append(sheet_text)

        wb.close()

        full_text = "\n\n".join(all_text_parts)
        full_text = self._strip_control_chars(full_text)

        metadata = {
            "format": "xlsx",
            "file_name": os.path.basename(path),
            "sheet_count": len(pages),
            "sheet_names": list(pages) if pages else [],
        }

        return self._make_result(full_text, pages, metadata)